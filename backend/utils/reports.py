import os
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pydantic import BaseModel, Field
from typing import List
from google.genai import Client
from google.genai import types

class SWOTData(BaseModel):
    strengths: List[str] = Field(description="List of strengths identified in the debate.")
    weaknesses: List[str] = Field(description="List of weaknesses identified in the debate.")
    opportunities: List[str] = Field(description="List of opportunities identified in the debate.")
    threats: List[str] = Field(description="List of threats identified in the debate.")

class RiskItem(BaseModel):
    description: str = Field(description="Description of the risk.")
    severity: str = Field(description="Severity (High, Medium, Low).")
    likelihood: str = Field(description="Likelihood (High, Medium, Low).")
    mitigation: str = Field(description="Proposed mitigation strategy.")

class FinancialsData(BaseModel):
    budget: str = Field(description="Total initial budget required (formatted as string, e.g. $150,000).")
    estimated_revenue: str = Field(description="Projected revenue (formatted as string, e.g. $450,000).")
    roi_3yr: str = Field(description="Projected 3-Year ROI (formatted as string, e.g. 150%).")
    cost_details: str = Field(description="Details on Capex, Opex, and cost structure.")

class TechReviewData(BaseModel):
    feasibility: str = Field(description="Feasibility rating (High, Medium, Low).")
    tech_stack: str = Field(description="Recommended technology stack (comma-separated list).")
    effort: str = Field(description="Estimated engineering effort (e.g. 4 months, 3 developers).")
    details: str = Field(description="Technical feasibility detail and details of recommendations.")

class ActionPhase(BaseModel):
    phase: str = Field(description="Name/number of the phase (e.g. Phase 1: Setup).")
    timeline: str = Field(description="Timeline for this phase (e.g. Month 1).")
    details: str = Field(description="Detailed tasks to perform in this phase.")

class StructuredReport(BaseModel):
    title: str = Field(description="Title of the strategic report.")
    executive_summary: str = Field(description="Executive summary of the debate, findings, and consensus (6-8 sentences).")
    swot: SWOTData = Field(description="SWOT analysis data.")
    risks: List[RiskItem] = Field(description="Matrix of identified compliance, operational, and financial risks.")
    financials: FinancialsData = Field(description="Financial KPI analysis details.")
    tech_review: TechReviewData = Field(description="Technical feasibility and stack review.")
    decision_score: int = Field(description="Overall strategic score from 1 (poor) to 100 (excellent).")
    recommendation: str = Field(description="Overall recommendation (e.g. Proceed with caution, approved with conditions, rejected).")
    action_plan: List[ActionPhase] = Field(description="Step-by-step action plan timeline phases.")

def compile_report_data(full_debate_text: str, proposal: str) -> dict:
    """Use Gemini's structured output mode to compile the full debate transcript into report data."""
    api_key = os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY")
    client = Client(api_key=api_key) if api_key else Client()
    
    prompt = f"""You are a senior business analyst compiling a formal strategic report.
Analyze the following corporate board debate transcript and extract the key findings into a structured report format.

Business Proposal: {proposal}

Board Debate Transcript:
{full_debate_text}
"""
    try:
        response = client.models.generate_content(
            model='gemini-1.5-flash',
            contents=prompt,
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
                response_schema=StructuredReport,
                temperature=0.1
            )
        )
        import json
        return json.loads(response.text)
    except Exception as e:
        print(f"Failed to compile report data using Gemini: {e}")
        return {
            "title": f"Strategic Report: {proposal[:30]}",
            "executive_summary": "Failed to extract debate summary automatically. Please review debate transcripts.",
            "swot": {"strengths": ["Clear objective"], "weaknesses": ["Adoption hurdles"], "opportunities": ["Market expansion"], "threats": ["Compliance risks"]},
            "risks": [{"description": "Data privacy", "severity": "Medium", "likelihood": "Medium", "mitigation": "Access control"}],
            "financials": {"budget": "$150,000", "estimated_revenue": "$300,000", "roi_3yr": "100%", "cost_details": "Opex and Capex details"},
            "tech_review": {"feasibility": "Medium", "tech_stack": "FastAPI, React", "effort": "4 Months", "details": "Feasible stack"},
            "decision_score": 70,
            "recommendation": "Approved with safety mitigations.",
            "action_plan": [{"phase": "Phase 1: Prototyping", "timeline": "Month 1", "details": "Core APIs setup"}]
        }


def generate_markdown_report(data: dict) -> str:
    """Generate a clean markdown report from board outcome data."""
    title = data.get("title", "Strategic Business Evaluation Report")
    summary = data.get("executive_summary", "No summary provided.")
    swot = data.get("swot", {"strengths": [], "weaknesses": [], "opportunities": [], "threats": []})
    risks = data.get("risks", [])
    financials = data.get("financials", {})
    tech_review = data.get("tech_review", {})
    decision_score = data.get("decision_score", 75)
    recommendation = data.get("recommendation", "Proceed with caution.")
    action_plan = data.get("action_plan", [])
    
    md = f"""# {title}

## Executive Summary
{summary}

## Decision Score: {decision_score}/100
**Overall Board Recommendation:** {recommendation}

## SWOT Analysis
| Category | Details |
|---|---|
| **Strengths** | {", ".join(swot.get("strengths", [])) or "None identified."} |
| **Weaknesses** | {", ".join(swot.get("weaknesses", [])) or "None identified."} |
| **Opportunities** | {", ".join(swot.get("opportunities", [])) or "None identified."} |
| **Threats** | {", ".join(swot.get("threats", [])) or "None identified."} |

## Risk Matrix
| Risk Description | Severity | Likelihood | Mitigation |
|---|---|---|---|
"""
    for risk in risks:
        md += f"| {risk.get('description', '')} | {risk.get('severity', '')} | {risk.get('likelihood', '')} | {risk.get('mitigation', '')} |\n"
        
    md += f"""
## Financial Projection & ROI
* **Initial Budget Required:** {financials.get('budget', 'TBD')}
* **Estimated Annual Revenue:** {financials.get('estimated_revenue', 'TBD')}
* **Projected 3-Year ROI:** {financials.get('roi_3yr', 'TBD')}
* **Cost Analysis Details:** {financials.get('cost_details', 'TBD')}

## Technical Feasibility Review
* **Feasibility Rating:** {tech_review.get('feasibility', 'Medium')}
* **Recommended Technology Stack:** {tech_review.get('tech_stack', 'TBD')}
* **Estimated Engineering Effort:** {tech_review.get('effort', 'TBD')}
* **Feasibility Details:** {tech_review.get('details', 'TBD')}

## Proposed Action Plan
"""
    for idx, action in enumerate(action_plan, 1):
        md += f"{idx}. **{action.get('phase', f'Phase {idx}')}**: {action.get('details', '')} (Timeline: {action.get('timeline', 'TBD')})\n"
        
    return md

def generate_pdf_report(data: dict, filepath: str) -> None:
    """Generate a high-quality, professional PDF report using ReportLab."""
    doc = SimpleDocTemplate(filepath, pagesize=letter, rightMargin=40, leftMargin=40, topMargin=40, bottomMargin=40)
    styles = getSampleStyleSheet()
    
    # Custom Styles
    title_style = ParagraphStyle(
        'ReportTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=24,
        leading=28,
        textColor=colors.HexColor('#0F172A'),
        spaceAfter=5
    )
    
    meta_style = ParagraphStyle(
        'MetaStyle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=11,
        leading=15,
        textColor=colors.HexColor('#0284C7'),
        spaceAfter=15
    )
    
    section_heading = ParagraphStyle(
        'SectionHeading',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=15,
        leading=18,
        textColor=colors.HexColor('#1E3A8A'),
        spaceBefore=14,
        spaceAfter=8,
        keepWithNext=True
    )
    
    body_style = ParagraphStyle(
        'ReportBody',
        parent=styles['BodyText'],
        fontName='Helvetica',
        fontSize=9.5,
        leading=13.5,
        textColor=colors.HexColor('#334155'),
        spaceAfter=6
    )

    header_style = ParagraphStyle(
        'TableHeader',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=9.5,
        leading=12,
        textColor=colors.white
    )

    story = []
    
    # Title & Metadata
    story.append(Paragraph(data.get("title", "Strategic Business Evaluation Report"), title_style))
    story.append(Paragraph(f"Boardroom AI Strategic Evaluation Score: {data.get('decision_score', 75)}/100 | Recommendation: {data.get('recommendation', 'Proceed')}", meta_style))
    story.append(Spacer(1, 10))
    
    # Executive Summary
    story.append(Paragraph("Executive Summary", section_heading))
    story.append(Paragraph(data.get("executive_summary", "No summary provided."), body_style))
    story.append(Spacer(1, 8))
    
    # SWOT
    story.append(Paragraph("SWOT Analysis", section_heading))
    swot = data.get("swot", {})
    swot_data = [
        [Paragraph("<b>Category</b>", header_style), Paragraph("<b>Identified Points</b>", header_style)],
        [Paragraph("<b>Strengths</b>", body_style), Paragraph(", ".join(swot.get("strengths", [])) or "None.", body_style)],
        [Paragraph("<b>Weaknesses</b>", body_style), Paragraph(", ".join(swot.get("weaknesses", [])) or "None.", body_style)],
        [Paragraph("<b>Opportunities</b>", body_style), Paragraph(", ".join(swot.get("opportunities", [])) or "None.", body_style)],
        [Paragraph("<b>Threats</b>", body_style), Paragraph(", ".join(swot.get("threats", [])) or "None.", body_style)]
    ]
    t_swot = Table(swot_data, colWidths=[110, 420])
    t_swot.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1E3A8A')),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#CBD5E1')),
        ('BACKGROUND', (0,1), (0,-1), colors.HexColor('#F8FAFC')),
        ('PADDING', (0,0), (-1,-1), 6),
    ]))
    story.append(t_swot)
    story.append(Spacer(1, 10))
    
    # Risks
    story.append(Paragraph("Risk Assessment Matrix", section_heading))
    risk_table_data = [
        [
            Paragraph("<b>Risk Description</b>", header_style),
            Paragraph("<b>Severity</b>", header_style),
            Paragraph("<b>Likelihood</b>", header_style),
            Paragraph("<b>Mitigation Strategy</b>", header_style)
        ]
    ]
    for r in data.get("risks", []):
        risk_table_data.append([
            Paragraph(r.get("description", ""), body_style),
            Paragraph(r.get("severity", ""), body_style),
            Paragraph(r.get("likelihood", ""), body_style),
            Paragraph(r.get("mitigation", ""), body_style)
        ])
    t_risk = Table(risk_table_data, colWidths=[150, 60, 60, 260])
    t_risk_styles = [
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1E3A8A')),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#CBD5E1')),
        ('PADDING', (0,0), (-1,-1), 6),
    ]
    # Alternating row background
    for idx in range(1, len(risk_table_data)):
        if idx % 2 == 0:
            t_risk_styles.append(('BACKGROUND', (0, idx), (-1, idx), colors.HexColor('#F8FAFC')))
    t_risk.setStyle(TableStyle(t_risk_styles))
    story.append(t_risk)
    story.append(Spacer(1, 10))
    
    # Financials
    story.append(Paragraph("Financial Feasibility Projections", section_heading))
    fin = data.get("financials", {})
    fin_text = f"<b>Initial Capex / Budget Required:</b> {fin.get('budget', 'TBD')}<br/>" \
               f"<b>Estimated 3-Year Revenues:</b> {fin.get('estimated_revenue', 'TBD')}<br/>" \
               f"<b>Projected 3-Year ROI:</b> {fin.get('roi_3yr', 'TBD')}<br/>" \
               f"<b>Details:</b> {fin.get('cost_details', 'TBD')}"
    story.append(Paragraph(fin_text, body_style))
    story.append(Spacer(1, 8))
    
    # Technical Review
    story.append(Paragraph("Technology Stack & Feasibility Review", section_heading))
    tech = data.get("tech_review", {})
    tech_text = f"<b>Feasibility Rating:</b> {tech.get('feasibility', 'Medium')}<br/>" \
                f"<b>Recommended Stack:</b> {tech.get('tech_stack', 'TBD')}<br/>" \
                f"<b>Engineering Effort Estimate:</b> {tech.get('effort', 'TBD')}<br/>" \
                f"<b>Details:</b> {tech.get('details', 'TBD')}"
    story.append(Paragraph(tech_text, body_style))
    story.append(Spacer(1, 10))
    
    # Action Plan
    story.append(Paragraph("Execution Timeline & Action Plan", section_heading))
    for idx, act in enumerate(data.get("action_plan", []), 1):
        act_text = f"<b>{idx}. {act.get('phase', f'Phase {idx}')}</b> (Timeline: {act.get('timeline', 'TBD')}): {act.get('details', '')}"
        story.append(Paragraph(act_text, body_style))
        
    doc.build(story)

def generate_pptx_report(data: dict, filepath: str) -> None:
    """Generate a high-quality PowerPoint Presentation summarizing the board's decision with clean layouts."""
    prs = Presentation()
    
    # Slide width & height (Standard 16:9 widescreen)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Cover Slide (Dark Slate Theme)
    slide_layout = prs.slide_layouts[5] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set dark background for cover slide
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42) # Slate-900
    
    # Title box
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", "Strategic Board Evaluation")
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.name = "Arial"
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Subtitle box
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(2))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = f"Prepared by BoardRoom AI\nRecommendation: {data.get('recommendation', 'Proceed')}\nScore: {data.get('decision_score', 75)}/100"
    p_sub.font.size = Pt(20)
    p_sub.font.name = "Arial"
    p_sub.font.color.rgb = RGBColor(148, 163, 184) # Slate-400
    
    # Helper to create styled light-themed slide
    def add_light_slide(title_text):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        bg = s.background
        fl = bg.fill
        fl.solid()
        fl.fore_color.rgb = RGBColor(248, 250, 252) # Slate-50
        
        # Add slide title
        tb = s.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1))
        tf_t = tb.text_frame
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.bold = True
        p_t.font.size = Pt(28)
        p_t.font.name = "Arial"
        p_t.font.color.rgb = RGBColor(30, 58, 138) # Navy
        return s
        
    # 1. Executive Summary Slide
    slide = add_light_slide("Executive Summary")
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("executive_summary", "No summary provided.")
    p.font.size = Pt(16)
    p.font.name = "Arial"
    p.font.color.rgb = RGBColor(51, 65, 85) # Slate-700
    
    # 2. SWOT Analysis Slide
    slide = add_light_slide("SWOT Analysis Matrix")
    swot = data.get("swot", {})
    categories = [
        ("Strengths", swot.get("strengths", []), Inches(0.75), Inches(1.8), RGBColor(16, 185, 129)),
        ("Weaknesses", swot.get("weaknesses", []), Inches(6.8), Inches(1.8), RGBColor(239, 68, 68)),
        ("Opportunities", swot.get("opportunities", []), Inches(0.75), Inches(4.5), RGBColor(59, 130, 246)),
        ("Threats", swot.get("threats", []), Inches(6.8), Inches(4.5), RGBColor(245, 158, 11))
    ]
    for title, items, left, top, color in categories:
        tb = slide.shapes.add_textbox(left, top, Inches(5.7), Inches(2.3))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.name = "Arial"
        p.font.color.rgb = color
        
        for it in items:
            p_it = tf.add_paragraph()
            p_it.text = f"  • {it}"
            p_it.font.size = Pt(13)
            p_it.font.name = "Arial"
            p_it.font.color.rgb = RGBColor(71, 85, 105)
            
    # 3. Financials & Tech Slide
    slide = add_light_slide("Financials & Technical Feasibility")
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Financial Feasibility:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(16, 185, 129)
    
    fin = data.get("financials", {})
    for k, v in [("Budget Required", fin.get("budget")), ("Revenues (3-Yr)", fin.get("estimated_revenue")), ("ROI (3-Yr)", fin.get("roi_3yr")), ("Details", fin.get("cost_details"))]:
        p = tf.add_paragraph()
        p.text = f"• {k}: {v}"
        p.font.size = Pt(13)
        p.font.name = "Arial"
        
    tx_box_tech = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.5), Inches(5))
    tf_tech = tx_box_tech.text_frame
    tf_tech.word_wrap = True
    p_tech = tf_tech.paragraphs[0]
    p_tech.text = "Technical Stack & Feasibility:"
    p_tech.font.bold = True
    p_tech.font.size = Pt(18)
    p_tech.font.color.rgb = RGBColor(59, 130, 246)
    
    tech = data.get("tech_review", {})
    for k, v in [("Feasibility Rating", tech.get("feasibility")), ("Technology Stack", tech.get("tech_stack")), ("Engineering Effort", tech.get("effort")), ("Feasibility Details", tech.get("details"))]:
        p = tf_tech.add_paragraph()
        p.text = f"• {k}: {v}"
        p.font.size = Pt(13)
        p.font.name = "Arial"
        
    # 4. Timeline & Action Slide
    slide = add_light_slide("Proposed Execution Timeline")
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = f"Strategic Recommendation: {data.get('recommendation')}"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(30, 58, 138)
    
    p = tf.add_paragraph()
    p.text = "Proposed Action Plan Phases:\n"
    p.font.bold = True
    p.font.size = Pt(16)
    
    for act in data.get("action_plan", []):
        p = tf.add_paragraph()
        p.text = f"  • {act.get('phase')} ({act.get('timeline')}): {act.get('details')}"
        p.font.size = Pt(14)
        p.font.name = "Arial"
        
    prs.save(filepath)
